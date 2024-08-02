# -*- coding: utf-8 -*-
"""
Created on Fri Aug  2 16:48:24 2024

@author: benny.ng
"""


import os
import comtypes.client
from pptx import Presentation

# Define os caminhos da pasta
pasta_pptx = ''

pptx_folder = os.path.join(os.path.dirname(__file__), pasta_pptx )
pdf_folder = os.path.join(pptx_folder, 'pdf')



def convert_pptx_to_pdf(pptx_folder, pdf_folder):
    # Cria a pasta de destino se ela não existir
    if not os.path.exists(pdf_folder):
        os.makedirs(pdf_folder)

    # Inicializa o PowerPoint
    powerpoint = comtypes.client.CreateObject("Powerpoint.Application")
    powerpoint.Visible = 1

    # Itera sobre todos os arquivos na pasta PPTX
    for filename in os.listdir(pptx_folder):
        if filename.endswith(".pptx"):
            pptx_path = os.path.join(pptx_folder, filename)
            pdf_path = os.path.join(pdf_folder, filename.replace(".pptx", ".pdf"))

            # Abre a apresentação
            presentation = powerpoint.Presentations.Open(pptx_path)
            # Exporta para PDF
            presentation.SaveAs(pdf_path, FileFormat=32)  # 32 é o formato PDF
            presentation.Close()

    # Fecha o PowerPoint
    powerpoint.Quit()
    
def renumerate_pptx(pptx_folder):
    for filename in os.listdir(pptx_folder):
        if filename.endswith(".pptx"):
            pptx_path = os.path.join(pptx_folder, filename)
            prs = Presentation(pptx_path)
            # slide_number
            for slide_idx in range(len(prs.slides)):
                slide = prs.slides[slide_idx]
                try:
                    name_shape = [
                        shape for shape in slide.shapes if shape.name == "slide_number"][0]
                    name_shape.text_frame.paragraphs[0].runs[0].text = str(slide_idx+1)
                except:
                    pass
                
            prs.save(pptx_path)
    
# Renumera slides
renumerate_pptx(pptx_folder)
# Converte os arquivos
convert_pptx_to_pdf(pptx_folder, pdf_folder)